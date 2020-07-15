#!/usr/bin/python

import pandas
import os
import glob
import json
import datetime
import argparse
from data_utils import *

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as pl
import matplotlib.dates as mdates


ta_dict = {'year': 'year',
           'conf': 'date',
           'five': '5 year period',
           'decade': 'decade'}
nm_dict = {'words': 'uses per million words',
           'conf': 'uses per conference'}

parser = argparse.ArgumentParser()
parser.add_argument('--output-dir', dest='output_dir', type=str, default='/mnt/d/Quentin/sunstone',
                    help='path for saving output files')
parser.add_argument('--search-data', dest='search_data', type=str, default='search_data.json',
                    help='JSON file containing search data')
parser.add_argument('--time-axis', dest='time_axis', type=str, default='year',
                    choices=ta_dict.keys(), help='specifies spacing of time axis')
parser.add_argument('--norm', dest='norm', type=str, default='words',
                    choices=nm_dict.keys(), help='specifies normalization factor')
parser.add_argument('--smooth', dest='smooth', type=int, default=None,
                    help='size of window for smoothing')
args = parser.parse_args()

# This needs to be set high if a large number of plots are being generated
pl.style.use('bmh')
prop_cycle = pl.cycler(color=['#e41a1c','#377eb8','#4daf4a','#984ea3',
                              '#ff7f00','#ffff33','#a65628','#f781bf','#999999'])
pl.rcParams.update({
    'figure.max_open_warning': 100,
    'lines.linewidth': 4,
    'axes.prop_cycle': prop_cycle,
    'font.size': 20,
    'figure.figsize': [10, 6],
    'figure.subplot.left': 0.125,
    'figure.subplot.right': 0.95,
    'figure.subplot.bottom': 0.15,
    'figure.subplot.top': 0.95})


# Setup PowerPoint presentation
from pptx import Presentation
from pptx.util import Inches, Pt
pres = Presentation()
title_slide_layout = pres.slide_layouts[0]
two_content_layout = pres.slide_layouts[3]
title_only_layout = pres.slide_layouts[5]
#title_only_layout.shapes[2].text_frame.paragraphs[0].font.size = Pt(40)
blank_layout = pres.slide_layouts[6]

slide = pres.slides.add_slide(title_slide_layout)
slide.shapes.title.text = 'Sunstone Presentation'
slide.placeholders[1].text = 'Quentin Spencer'

def create_ppt_slide(fname, title_str):
    '''Add figure slide to PowerPoint presentation'''
    full_filename = '{0:s}/{1:s}.png'.format(args.output_dir, fname)
    pl.savefig(full_filename)
    slide = pres.slides.add_slide(title_only_layout)
    slide.shapes.title.text = title_str
    pic = slide.shapes.add_picture(full_filename,
                                   Inches(0), Inches(1.5), height=Inches(6))
    # Move picture behind title
    cursor_sp = slide.shapes[0]._element
    cursor_sp.addprevious(pic._element)
    print('added {0:s} to presentation'.format(full_filename))


print("Loading data")
searches = json.load(open(args.search_data))


decade_word_counts = pandas.read_table('data/general_counts.txt')
decade_word_counts = decade_word_counts.set_index('decade')
scripture_word_counts = {'Bible (KJV)': 783137,
                         'Book of Mormon': 267846,
                         'Doctrine and Covenants': 115658,
                         'Pearl of Great Price': 28366}
scripture_data = pandas.Series(scripture_word_counts).to_frame('words')
scripture_data_r = (scripture_data/1000).round().astype(int)*1000


fig, ax = pl.subplots()
ax.set_prop_cycle(prop_cycle[1:])
labels = scripture_data_r.index + '\n' + \
         scripture_data_r['words'].map(lambda x: '{:,d}'.format(x))
ax.pie(scripture_data_r['words'],
       labels=labels,
       labeldistance=1.25,
       wedgeprops={'linewidth':2, 'edgecolor':'white'},
       startangle=35,
       counterclock=False)
pl.subplots_adjust(bottom=0.05, top=0.9, right=0.85)
create_ppt_slide('std_works',
                 'The LDS standard works\n({0:,d} words)'.format(
                     scripture_data_r.sum()[0]))


fig, ax = pl.subplots()
mean_words = decade_word_counts['words'].mean()/1e6
dwc_plot_data = decade_word_counts.drop('talks', 1).append(scripture_data.transpose())
dwc_plot_data = dwc_plot_data[-1:].append(dwc_plot_data[:-1]).fillna(0)/1e6
dwc_plot_data.plot.bar(ax=ax, stacked=True, legend=False, width=0.8)
pl.axhline(y=mean_words,
           color='black', linewidth=1, zorder=0)
pl.text(len(dwc_plot_data)-0.1, mean_words, '{0:.2f}'.format(mean_words),
        horizontalalignment='left', verticalalignment='center')
labels = [''] + list(decade_word_counts.index)
pl.xlabel('decade')
pl.xticks(range(len(labels)), labels, rotation=45, ha='right')
pl.ylabel('million words')
pl.subplots_adjust(bottom=0.25, right=0.92)
pl.grid(axis='x')
create_ppt_slide('decades',
                 'Total words spoken in general conference per decade')


fig, ax = pl.subplots()
talk_length = decade_word_counts['words']/decade_word_counts['talks']
talk_length.plot.bar(ax=ax, width=0.8)
pl.xlabel('decade')
pl.xticks(range(len(talk_length)), talk_length.index, rotation=45, ha='right')
pl.ylabel('words')
pl.subplots_adjust(bottom=0.25)
pl.grid(axis='x')
create_ppt_slide('talk_length',
                 'General conference talk length')

def bar_chart(df, legend=True):
    fig, ax = pl.subplots()
    df.plot.bar(ax=ax, width=0.8, stacked=True)
    pl.xticks(range(len(df)), df.index, rotation=45, ha='right')    
    pl.subplots_adjust(bottom=0.30)
    pl.grid(axis='x')
    if legend:
        pl.legend(df.columns)

#####################################
# Load texts & reference data
df_all = load_data(source='file')
talks_only = get_only_talks(df_all)

apostle_data = load_apostle_data()
apostle_data['short name'] = apostle_data['name'].map(shorten_name)
apostle_sn_dict = apostle_data[['name', 'short name']].set_index('name').to_dict()['short name']
president_list = apostle_data[~apostle_data['sdate_p'].isna()]['name'].to_list()
recent_presidents = talks_only['president'].unique()
recent_presidents_ln = [x.split(' ')[-1] for x in recent_presidents]

#####################################

talk_count = talks_only.groupby('author').size().sort_values(
    ascending=False)[:20].to_frame('talks')
talk_count['pres'] = ~talk_count.index.isin(president_list)
talk_count = talk_count.pivot(columns='pres', values='talks').fillna(0)
talk_count.columns = ['Presidents', 'Apostles']
talk_count['sum'] = talk_count.sum(1)
talk_count = talk_count.sort_values('sum', ascending=False).drop('sum', 1)
talk_count.index = [apostle_sn_dict[x] for x in talk_count.index]

bar_chart(talk_count)
pl.subplots_adjust(bottom=0.32, left=0.13)
pl.ylabel('conference talks given')
create_ppt_slide('prolific_speakers',
                 'Most prolific speakers\n1942-present')

#####################################

talk_count_na = talks_only[talks_only['rank']==100].groupby('author').size()
talk_count_na = talk_count_na.to_frame('talks').sort_values('talks', ascending=False)[:15]
talk_count_na['title'] = 'Seventy'
talk_count_na.loc['Eldred G. Smith', 'title'] = 'Patriarch to the Church'
talk_count_na.loc[['Victor L. Brown', 'Joseph L. Wirthlin'], 'title'] = 'Presiding Bishop'
talk_count_na = talk_count_na.pivot(columns='title', values='talks').fillna(0)
talk_count_na['sum'] = talk_count_na.sum(1)
talk_count_na = talk_count_na.sort_values('sum', ascending=False).drop('sum', 1)

bar_chart(talk_count_na)
pl.gca().set_ylim(top=70)
pl.subplots_adjust(bottom=0.38, left=0.15)
pl.ylabel('conference talks given')
create_ppt_slide('non_apostles',
                 'Most prolific non-apostle speakers 1942-present')

#####################################

current_q15 = talks_only[(talks_only['date']==talks_only.iloc[-1]['date']) &
                         (talks_only['rank']<20)]['author'].drop_duplicates().to_list()
current_q15_talk_count = talks_only[(talks_only['author'].isin(current_q15)) &
                                    (talks_only['rank']>20)].groupby('author').size()
asd = apostle_data[apostle_data['dod'].isna()][['name', 'sdate']].set_index('name')
current_q15_talk_count = current_q15_talk_count.to_frame('talks').join(
    asd, how='outer').fillna(0).sort_values('sdate').drop('sdate', 1)
current_q15_talk_count.index = current_q15_talk_count.index.map(lambda x: x.split(' ')[-1])
bar_chart(current_q15_talk_count, legend=False)
pl.ylabel('conference talks given')
create_ppt_slide('pre_apostles',
                 'Talks given by current apostles prior to call')

#####################################

# Length of GC speaking careers
date_grp = talks_only.groupby('author')['date']
speaking_career_length = (date_grp.max() - date_grp.min()).apply(
    lambda x: round(x.days/365, 1)).sort_values(ascending=False)

#####################################
first_year = df_all.year.min()
last_year = df_all.year.max()
daterange = [first_year, last_year]
group = 'year'

# Process references
ref_df = get_scripture_refs(df_all)
ref_counts_raw = get_ref_counts(ref_df, group)
ref_rate = ref_counts_raw.divide(talks_only.groupby(group).sum()['word_count'], 0)*1e6
rpm_str = 'references per million words'
ref_rate.columns = ref_rate.columns.get_level_values(1)
ref_rate.columns.name = 'Standard Work'
ref_rate_total = ref_rate.sum(1)
ref_rate_frac = ref_rate.divide(ref_rate_total, 0)*100

ref_freq = ref_df.groupby('ref').count()[group].to_frame('uses')
num_refs = 20
top_refs = ref_freq.sort_values('uses').iloc[-num_refs:][::-1]

top_refs_sw = ref_df[['ref', 'sw']].drop_duplicates().merge(top_refs, left_on='ref', right_on='ref')
top_refs_swcols = top_refs_sw.pivot(index='ref', columns='sw', values='uses').fillna(0).astype(int)
top_refs_swcols['all'] = top_refs_swcols.sum(1)
top_refs_swcols = top_refs_swcols.sort_values('all', ascending=False).drop('all', 1)
top_refs_swcols.columns = pandas.Index(
    ['Book of Mormon', 'Doctrine & Covenants',
     'New Testament', 'Pearl of Great Price'], name='standard work')


#######################################
fig, ax = pl.subplots()
ref_rate_total.plot(ax=ax)
ax.set_xlim(daterange)
ax.set_ylim(bottom=0)
pl.ylabel(rpm_str)
pl.subplots_adjust(left=0.15)
create_ppt_slide('refs_total',
                 'Scripture citation rate')

#####################################
fig, ax = pl.subplots()
ref_rate.plot(ax=ax)
#ref_rate.rolling(5, min_periods=1, center=True).mean().plot(ax=ax)
ax.set_xlim(daterange)
ax.set_ylim(bottom=0)
pl.ylabel(rpm_str)
pl.legend(ref_rate.columns, ncol=2, loc='upper left')
create_ppt_slide('refs',
                 'Scripture citations by standard work')

#####################################
fig, ax = pl.subplots()
ref_rate_frac.plot.area(ax=ax, legend=False)
ax.set_xlim(daterange)
ax.set_ylim(top=100)
pl.legend(ref_rate_frac.columns, ncol=2, loc='upper left')
pl.ylabel('percentage of total')
create_ppt_slide('refs_frac',
                 'Scripture citations by standard work')

#######################################
fig, ax = pl.subplots()
top_refs_swcols.plot.bar(stacked=True, width=0.8)
#ax.set_xlim([-0.5, num_refs-0.5])
pl.xticks(range(num_refs), top_refs.index, rotation=45, ha='right')
pl.subplots_adjust(bottom=0.25)
pl.grid(axis='x')
pl.ylabel('total references')
create_ppt_slide('top_scriptures',
                 'Most cited scriptures')

#####################################

def barplot(axes, data, title):
    data[::-1].plot.barh(ax=axes, width=0.8)
    axes.grid(axis='y')
    axes.set_title(title)
    axes.set_xlabel('total uses')

def text_search_plot(search, filename, title_str, leg_loc=None, spacer=' '):
    if 'top_user' not in search.keys():
        search.update({'top_user': 'true'})
    result = text_search(talks_only, search, spacer=spacer)
    fig, ax = pl.subplots()
    result.plot(ax=ax, legend=False)
    pl.legend(result.columns, loc=leg_loc)
    pl.ylabel(rpm_str)
    pl.grid(axis='x')
    ax.set_xlim(daterange)
    ax.set_ylim(bottom=0)
    create_ppt_slide(filename, title_str)

text_search_plot({'search': [{'include': '(grief|grieve|grieving)',
                              'label':'grief/grieving'},
                             {'include': 'gratitude'}]},
                 'gratitude',
                 'The Sunstone theme')

fig, ax = pl.subplots(1, 2)
pl.subplots_adjust(left=0.22, wspace=0.8, top=0.88)
users1 = top_users(talks_only, '(grief|grieve|grieving)')['matches']
users1.index = [apostle_sn_dict[x] for x in users1.index]
users2 = top_users(talks_only, 'gratitude')['matches']
users2.index = [apostle_sn_dict[x] for x in users2.index]
barplot(ax[0], users1, 'grief')
barplot(ax[1], users2, 'gratitude')
create_ppt_slide('gratitude2', 'The Sunstone theme')

#####################################

top_ngrams = pandas.read_pickle('top_ngrams.pkl')

def term_bar_plot(dataframe, filename, title_str):
    cap_list = ['indians', 'jesus', 'christ', 'mormonism', 'i ']
    for cap_word in cap_list:
        dataframe.index = dataframe.index.str.replace(cap_word, cap_word.capitalize())

    fig, ax = pl.subplots()
    ng_plot_data = dataframe['ratio']
    max_len = ng_plot_data.index.str.len().max()
    left_ = 0.2 + (max_len - 10)/25*0.4
    ng_plot_data.index = [ldqm+x+rdqm for x in ng_plot_data.index]
    ng_plot_data[::-1].plot.barh(ax=ax, width=0.8)
    pl.subplots_adjust(left=left_, bottom=0.13)
    pl.grid(axis='y')
    pl.xlabel('max/min ratio')
    create_ppt_slide(filename, title_str)


term_bar_plot(top_ngrams[(top_ngrams['peak']<=1960) & (top_ngrams['ratio']>55)],
              'b1960',
              'most-changed usages,\npeak before 1970')
#'communism',-
#'constitution',-
#'indians',-
#'free agency',-
#'our boys',-
#'this land',-
#'fullness',
#'mormonism',-
#'this nation',-
#'fellow men'

for searchdata, filename, title_str in (
        ({'search': [{'include': 'our boys'}]},
         'ww2',
         'World War II phrases: ' + ldqm + 'our boys' + rdqm),
        ({'search': [{'include': 'communism'}]},
         'cold war',
         'Cold War phrases: ' + ldqm + 'communism' + rdqm),
        ({'search': [{'include': 'constitution'},
                     {'include': 'this land'},
                     {'include': 'this nation'}]},
         'america',
         'America-centric phrases'),
        ({'search': [{'include': 'indians', 'label': 'Indians'}]},
         'indians',
         ldqm + 'Indians' + rdqm),
        ):
    text_search_plot(searchdata, filename, title_str)

text_search_plot(
    {"search": [{"label": "free agency",
	         "include": "free agency"},
	        {"label": "moral agency",
	         "include": "moral agency"},
	        {"label": "other usages of \"agency\"",
	         "include": "agency",
	         "exclude": ["free agency", "moral agency"]}]},
    'agency',
    'evolution of the usage of '+ldqm+'agency'+rdqm,
    leg_loc='upper left')

for searchdata, filename, title_str in (
        ({'search': [{'include': 'mormonism', 'label': 'Mormonism'}]},
         'mormon',
         ldqm + 'Mormonism' + rdqm),
                ({'search': [{'include': 'fullness'},
                     {'include': 'fulness'}]},
         'fulness',
         'spelling variations'),
        ({'search': [{'include': 'fellow men'},
                     {'include': 'fellowmen'}]},
         'fellowmen',
         'style guide change?'),
        ({'search': [{'include': 'brethren and sisters'},
                     {'include': 'brothers and sisters'},
                     {'include': 'the brethren'}]},
         'brethren',
         'language modernization'),
        ):
    text_search_plot(searchdata, filename, title_str)

term_bar_plot(top_ngrams[(top_ngrams['peak'].isin([1970,1980])) & (top_ngrams['ratio']>55)],
              '1970-1980',
              'most-changed usages,\npeak 1970-1990')
#'welfare services',-
#'home teachers',-
#'family home evening',-
#'home teaching',-
#'meaningful',-
#'resource',-
#'communicate',
#'stewardship',-
#'our fellowmen',
#'self reliance',-
#'concern for',
#'regional',-
#'priesthood leaders',
#'home teacher',-
#'drugs'-

for searchdata, filename, title_str in (
        ({'search': [{'include': 'welfare services'},
                     {'include': 'resource'},
                     {'include': 'self.reliance', 'label': 'self-reliance'}]},
         'welfare',
         'Welfare Session'),
        ({'search': [{'include': 'drugs'}]},
         'drugs',
         ldqm + 'drugs' + rdqm),
        ({'search': [{'include': 'regional'}]},
         'regional',
         ldqm + 'regional' + rdqm + ': church organizational structures'),
        ({'search': [{'include': 'home (teacher|teaching)', 'label': 'home teacher/teaching'},
                     {'include': 'family home evening'}]},
         'programs',
         'New church programs'),
        ({'search': [{'include': 'stewardship'},
                     {'include': 'meaningful'}]},
         'other',
         'other phrases'),
        ):
    text_search_plot(searchdata, filename, title_str, spacer='\n')


term_bar_plot(top_ngrams[(top_ngrams['peak']>=1990) & (top_ngrams['ratio']>95)],
              'a1990',
              'most-changed usages,\npeak after 1990')
#'commitment',-
#'the atonement of jesus christ',-
#'discipleship',-
#'pornography',-
#'priesthood holders',
#'family history',-
#'plan of happiness',
#'challenges',-
#'i invite',
#'focus on',-
#'media',-
#'scripture study',-
#'internet',
#'focused',-
#'priesthood authority',-
#'in the sacred name of jesus christ'

#####################################

text_search_plot({'search': [{'include': 'commitment'},
                             {'include': 'challenges'},
                             {'include': 'focus'}]},
                 'commitment',
                 'General linguistic trends\nin the late 20th century')

#####################################

google_data = pandas.read_csv('google_book_search.csv')
google_data = google_data.set_index(google_data['year'].map(
    lambda x: pandas.to_datetime(datetime.date(x, 1, 1), utc=True))).drop('year', 1)*1e6
# drop last rows where data looks inconsistent
google_data = google_data[:-4]

fig, ax = pl.subplots()
google_data.plot(ax=ax)
pl.ylabel(rpm_str)
pl.grid(axis='x')
ax.set_xlim(daterange)
ax.set_ylim(bottom=0)
create_ppt_slide('google', 'Google book search results')

#####################################

for searchdata, filename, title_str in (
        ({'search': [{'include': 'genealogy'},
                     {'include': 'family history'}]},
         'fam_hist',
         'renaming genealogy'),
        ({'search': [{'include': 'plan of happiness'},
                     {'include': 'plan of salvation'}]},
         'plan_happiness',
         ldqm + 'plan of happiness' + rdqm),
        ({'search': [{'include': 'phone'},
                     {'include': 'media'},
                     {'include': 'internet'}]},
         'tech',
         'New technologies'),
        ({'search': [{'include': 'i invite you', 'label': 'I invite you'}]},
         'invite',
         'asking more politely: ' + ldqm + 'I invite you' + rdqm),
        ({'search': [{'include': 'pornography'}]},
         'porn',
         'worries about pornography'),
        ({'search': [{'include': 'scripture study'}]},
         'scripturestudy',
         ldqm + 'scripture study' + rdqm),
):
    text_search_plot(searchdata, filename, title_str)

for searchdata, filename, title_str in (
        ({"search": [{"label": "daily prayer",
	              "include": "daily.{,30}prayer"},
	             {"label": "scripture study/reading",
	              "include": "(scripture study|scripture reading|study[^.,]{,15}scripture)"},
	             {"label": "daily scripture study/reading",
	              "include": "(scriptures every day|daily[^.,-]{,30}scripture|daily[^.,-]{,30}[^b]read)"}]},
         'practices',
         'an emphasis on daily practices'),
        ({'search': [{'include': 'the atonement of jesus christ',
                      'label': 'the atonement of Jesus Christ'},
                     {'include': 'discipleship'}]},
         'jesus',
         'A new emphasis on atonement & Jesus'),
        ({'search': [{'include': 'priesthood leader'},
                     {'include': 'priesthood holder'},
                     {'include': 'priesthood authority'}]},
         'priesthood',
         'priesthood & authority'),
        ):
    text_search_plot(searchdata, filename, title_str, spacer='\n')

#####################################

search = {'search': [{'include': 'priesthood leader'},
                     {'include': 'priesthood holder'},
                     {'include': 'priesthood power'},
                     {'include': 'priesthood authority'},
                     {'include': 'priesthood keys'}]}
result = text_search(talks_only, search)
fig, ax = pl.subplots()
result.plot.area(ax=ax)
pl.ylabel(rpm_str)
pl.grid(axis='x')
ax.set_xlim(daterange)
ax.set_ylim(bottom=0)
create_ppt_slide('priesthood2', 'priesthood & authority')
#pl.legend(ref_rate_frac.columns, ncol=2, loc='center left')

text_search_plot({'search': [{'include': 'in the sacred name of jesus christ'}]},
                 'sacrednameofjesus',
                 ldqm + '...in the sacred name of Jesus Christ' + rdqm,
                 spacer='\n')

search_data_jesus = {
    "case sensitive": "true",
    #"include sum": "true",
    "search": [{"label": "Jesus Christ",
	        "include": "Jesus Christ",
	        "exclude": ["[Cc]hurch of Jesus Christ",
			    "Jesus Christ.{0,20} [Aa]men"]},
	       {"label": "the Savior",
	        "include": "[Tt]he Savior"},
	       {"label": "Jesus",
	        "include": "Jesus(?! Christ)"},
	       {"label": "Christ",
	        "include": "(?<!Jesus )Christ",
	        "exclude": ["Jesus is the Christ"]},
	       {"label": "Master",
	        "include": "Master"},
	       {"label": "Redeemer",
	        "include": "Redeemer"}]}
text_search_plot(search_data_jesus, 'namesforjesus',
                 'Names for Jesus', leg_loc='upper left')
result_j = text_search(talks_only, search_data_jesus)
result_j_total = result_j.sum(1)
result_j_norm = result_j.divide(result_j_total, 0)*100

fig, ax = pl.subplots()
result_j_norm.plot.area(ax=ax, legend=False)
pl.legend(result_j_norm.columns.map(lambda x: x.split(' [')[0]),
          loc='upper left', ncol=2)
pl.ylabel('% of references')
pl.grid(axis='x')
ax.set_xlim(daterange)
ax.set_ylim(bottom=0, top=100)
create_ppt_slide('namesforjesus2', 'Names for Jesus')


search_data_satan = {#"include sum": "true",
    "search": [{"label": "the Adversary", "include": "the adversary"},
               {"label": "Satan",   "include": "satan"},
	       {"label": "Lucifer", "include": "lucifer"},
	       {"label": "the Devil", "include": "the devil"}]}
text_search_plot(search_data_satan, 'namesforsatan',
                 'Names for Satan', leg_loc='upper left')
result_s = text_search(talks_only, search_data_satan)
result_s_total = result_s.sum(1)
result_s_norm = result_s.divide(result_s_total, 0)*100

fig, ax = pl.subplots()
result_s_norm.plot.area(ax=ax, legend=False)
pl.legend(result_s_norm.columns.map(lambda x: x.split(' [')[0]),
          loc='upper left', ncol=2)
pl.ylabel('% of references')
pl.grid(axis='x')
ax.set_xlim(daterange)
ax.set_ylim(bottom=0, top=100)
create_ppt_slide('namesforsatan2', 'Names for Satan')


js_compare = result_j_total.to_frame('Jesus').join(result_s_total.to_frame('Satan'))
fig, ax = pl.subplots()
js_compare.plot(ax=ax)
pl.ylabel(rpm_str)
pl.grid(axis='x')
ax.set_xlim(daterange)
ax.set_ylim(bottom=0)
create_ppt_slide('jesusvsatan', 'References to Jesus & Satan')


#####################################
sa_data = {'grace': ['grace', 'mercy', 'mercies', 'merciful'],
           'works': ['obey', 'obedient', 'qualify', 'qualified', 'worthy', 'worthiness']}
speaker_averages = get_speaker_refs(talks_only, data=sa_data)
speaker_averages = speaker_averages.rename(columns={'President': 'church president'})
#####################################

president_list_t = talks_only['president'].unique().tolist()
apostle_list = talks_only[(talks_only['date']==talks_only.iloc[-1]['date']) & \
                          (talks_only['rank']>1) & (talks_only['rank']<20)]['author'].unique().tolist()
dead_apostle_list = ['Boyd K. Packer', 'Neal A. Maxwell', 'Bruce R. McConkie', 'James E. Faust']

q15_averages = speaker_averages[speaker_averages.index.isin(apostle_list)]
pres_averages = speaker_averages[speaker_averages.index.isin(president_list_t)]
prev_q15_averages = speaker_averages[speaker_averages.index.isin(dead_apostle_list)]

note_data = q15_averages.append(pres_averages).append(prev_q15_averages).assign(ha='right', va='bottom')
note_data = note_data.assign(last_name=note_data.index.map(lambda x: x.split(' ')[-1]))
note_data.loc['George Albert Smith', 'last_name'] = 'GASmith'
note_data = note_data.reset_index().set_index('last_name')

def create_scatter_fig(col0, col1, note_data, titlestr, filename):

    def annotate_group(n_data):
        ret_list = []
        for ln, data in n_data.iterrows():
            if data.name=='George Albert Smith':
                name = 'GASmith'
            elif data.name=='Joseph Fielding Smith':
                name = 'JFSmith'
            else:
                name = data.name.split(' ')[-1]
            an = ax.annotate(name, data[[col0, col1]], rotation=data['r'],
                             ha=data['ha'], va=data['va'], size='small')
            ret_list.append(an)
        
    nd = note_data.reset_index().set_index('author')

    fig, ax = pl.subplots()
    msize = 80
    colors = prop_cycle.by_key()['color']
    nd.plot.scatter(col0, col1, ax=ax, s=0)
    ax.set_autoscale_on(False)
    pl.plot([0,3], [0,3], c=[0.7,0.7,0.7])
    ax.set_xlim(left=0)
    ax.set_ylim(bottom=0)
    
    pres_averages.plot.scatter(col0, col1, ax=ax, c=colors[1], s=msize)
    annotate_group(nd.loc[pres_averages.index.to_list()])
    create_ppt_slide(filename+'1', titlestr)

    prev_q15_averages.plot.scatter(col0, col1, ax=ax, c=colors[2], s=msize)
    annotate_group(nd.loc[prev_q15_averages.index.to_list()])    
    create_ppt_slide(filename+'2', titlestr)

    q15_averages.plot.scatter(col0, col1, ax=ax, c=colors[0], s=msize)
    annotate_group(nd.loc[q15_averages.index.to_list()])    
    create_ppt_slide(filename+'3', titlestr)

#############################

note_data['ha'] = 'right'
note_data['va'] = 'bottom'
note_data['r'] = 0
note_data.loc[
    ['Renlund', 'Monson', 'Cook', 'Bednar', 'Faust', 'Nelson', 'Eyring',
     'Hunter', 'Smith', 'Stevenson', 'McConkie', 'Oaks', 'Holland',
     'Uchtdorf', 'Andersen'], 'ha'] = 'left'
note_data.loc[
    ['Christofferson', 'Cook', 'Oaks', 'Hinckley', 'Faust', 'McKay', 'Benson',
     'Holland', 'Andersen', 'Soares', 'Monson', 'Ballard', 'Kimball'], 'va'] = 'top'
titlestr = 'Jesus vs. Satan references\n(per 1000 words)'
create_scatter_fig('Jesus', 'Satan', note_data, titlestr, 'jesus_vs_satan')




#######################################
fig, ax = pl.subplots()
pres_df = talks_only[['date', 'year', 'decade']].copy()

for pr in president_list:
    idx = talks_only['president']==pr
    lastname = pr.split(' ')[-1]
    if pr=='Joseph Smith':
        altstr = 'Prophet Joseph(?! Smith)'
    elif lastname=='Smith': # "President Smith" is ambiguous
        altstr = 'President Smith'
    else:
        altstr = 'President ' + lastname
    pres_df[pr] = talks_only['body'].str.count(pr) + \
                  talks_only['body'].str.count(altstr)
    pres_df.loc[idx, 'current'] = talks_only.loc[idx, 'body'].str.count(pr) + \
                                  talks_only.loc[idx, 'body'].str.count(altstr)
pres_refs = pres_df.groupby('year').sum()
pres_refs = pres_refs.divide(talks_only.groupby(group).sum()['word_count'], 0)*1e6
pres_refs[recent_presidents].plot(ax=ax)
ax.set_xlim(daterange)
ax.set_ylim(bottom=0)
pl.legend(recent_presidents_ln, ncol=3, loc='upper center')
pl.ylabel(rpm_str)
create_ppt_slide('presidents',
                 'References to\npresidents of the church')


fig, ax = pl.subplots()
pres_refs['current'].plot(ax=ax)
ax.set_xlim(daterange)
ax.set_ylim(bottom=0)
pl.ylabel(rpm_str)
#pl.ylabel('references per conference')
create_ppt_slide('cur_pres',
                 'References to the current president of the church')

##############################

note_data['ha'] = 'left'
note_data['va'] = 'bottom'
note_data['r'] = 0
note_data.loc[['Hinckley', 'Gong', 'Cook', 'McKay', 'Lee', 'Soares', 'Kimball',
               'Benson', 'Packer', 'Ballard', 'GASmith', 'Oaks'], 'ha'] = 'right'
note_data.loc[['Maxwell', 'McConkie'], 'ha'] = 'center'
note_data.loc[['Cook', 'Stevenson', 'Andersen', 'Rasband', 'Renlund', 'Monson', 'Hunter',
               'Ballard', 'Maxwell', 'Kimball', 'Hinckley'], 'va'] = 'top'
note_data.loc[['Eyring', 'Holland'], 'r'] = 15
note_data.loc[['Christofferson', 'Bednar'], 'r'] = 10
titlestr = 'Jesus vs. church president\n(per 1000 words)'
create_scatter_fig('Jesus', 'church president', note_data, titlestr, 'jesus_vs_pres')

##############################
pres_cites = []
for pr in president_list:
    pres_data = apostle_data[apostle_data['name']==pr].iloc[0]
    pres_idx = pres_df.date > pres_data['edate_p'] + datetime.timedelta(180)
    pres_avg = 0 if pres_idx.sum()==0 else \
               pres_df[pres_idx][pr].sum()/talks_only[pres_idx]['word_count'].sum()*1e6
    pres_cites.append([pr, pres_avg, pres_data['edate_p']])
post_cites_pres = pandas.DataFrame(pres_cites, columns=('pres', 'cites', 'death')).set_index('pres')
post_cites_pres = post_cites_pres[post_cites_pres['cites']>0].sort_values('cites', ascending=False)



fig, ax = pl.subplots()
#ax.bar(range(len(post_cites_pres)), post_cites_pres['cites'], align='center')
#ax.set_xlim([-0.5, len(post_cites_pres)-0.5])
post_cites_pres['cites'].plot.bar(ax=ax, rot=45, width=0.8)
pl.xticks(range(len(post_cites_pres)),
          [apostle_sn_dict[x] for x in post_cites_pres.index],
          rotation=45, ha='right')
pl.subplots_adjust(bottom=0.30)
pl.grid(axis='x')
pl.ylabel(rpm_str)
create_ppt_slide('top_pres',
                 'Most mentioned presidents after death')


#################################################

note_data['ha'] = 'left'
note_data['va'] = 'bottom'
note_data['r'] = 0
note_data.loc[['GASmith', 'Gong', 'Benson', 'Grant', 'McKay', 'Holland',
               'Soares', 'Hinckley', 'Kimball', 'Ballard', 'Eyring',
               'Faust', 'Maxwell', 'Monson'], 'ha'] = 'right'
note_data.loc[['Uchtdorf', 'Monson', 'Holland', 'Lee'], 'ha'] = 'center'
note_data.loc[['Grant'], 'va'] = 'center'
note_data.loc[['McConkie', 'Holland', 'Andersen', 'Monson'], 'va'] = 'top'
note_data.loc[['Packer', 'Rasband', 'Oaks'], 'r'] = 15
titlestr = 'Jesus vs. Joseph\n(per 1000 words)'
create_scatter_fig('Jesus', 'Joseph', note_data, titlestr, 'jesus_vs_joseph')

#################################################
pres_cites_all = []
recent_pres_list = talks_only['president'].unique().tolist()
for pr in recent_pres_list[:-1]:
    pres_data = apostle_data[apostle_data['name']==pr].iloc[0]
    confs_during = (pres_refs.index <= pres_data['edate_p']) & (pres_refs.index >= pres_data['sdate_p'])
    confs_post = pres_refs.index > pres_data['edate_p'] + datetime.timedelta(180)
    nc_post = sum(confs_post)
    pres_cites_all.append(
        [pr,
         sum(pres_refs[confs_during][pr])/float(sum(confs_during)),
         sum(pres_refs[confs_post][pr])/float(nc_post) if nc_post else 0])
ps_all = pandas.DataFrame(
    pres_cites_all,
    columns=('pres', 'cites_during', 'cites_post')).set_index('pres')
#################################################


all_q15 = apostle_data['name'].to_list()
all_q12 = list(set(all_q15) - set(president_list))

apo_counts = []
for apo in all_q12:
    citations = talks_only.body.str.count(apo)
    apo_data = apostle_data[apostle_data['name']==apo].iloc[0]
    all_idx = talks_only.date > apo_data['sdate']
    post_idx = talks_only.date > apo_data['edate'] + datetime.timedelta(180)
    all_ratio = 0 if all_idx.sum()==0 else \
                citations[all_idx].sum()/talks_only[all_idx]['word_count'].sum()*1e6
    post_ratio = 0 if post_idx.sum()==0 else \
                 citations[post_idx].sum()/talks_only[post_idx]['word_count'].sum()*1e6
    apo_counts.append([apo, all_ratio, post_ratio, apo_data['edate']])

apo_cites = pandas.DataFrame(apo_counts, columns=('name', 'all', 'post', 'death')).fillna(0)
append_rows = apo_cites[apo_cites.post>11.4].set_index('name')[['post', 'death']]
post_cites_pres.columns = ('Presidents', 'death')
append_rows.columns = ('Apostles', 'death')
all_cites = post_cites_pres.append(append_rows, sort=False).fillna(0)
all_cites = all_cites.drop('Joseph Smith')
all_cites.index = [apostle_sn_dict[x] for x in all_cites.index]
all_cites['total'] = all_cites.sum(1)
ac_final = all_cites.sort_values('total', ascending=False).drop(['total', 'death'], 1)

bar_chart(ac_final)
pl.ylabel(rpm_str)
create_ppt_slide('influencers',
                 'Posthumous mentions of past leaders')

ac_chron = all_cites.sort_values('death').drop(['total', 'death'], 1)
bar_chart(ac_chron)
pl.ylabel(rpm_str)
create_ppt_slide('influencers_chron',
                 'Posthumous mentions of past leaders (chronological by death date)')

##############################

slide = pres.slides.add_slide(two_content_layout)
slide.shapes.title.text = 'Grace & Works'
slide.placeholders[1].text = '\n'.join(sa_data['grace'])
slide.placeholders[2].text = '\n'.join(sa_data['works'])

search_data_gw = {
    'top_user': 'false',
    'search': [{'include': '('+'|'.join(sa_data['grace'])+')',
                'label': 'grace-oriented words'},
               {'include': '('+'|'.join(sa_data['works'])+')',
                'label': 'works-oriented words'}]}
text_search_plot(search_data_gw, 'graceandworks', 'grace & works')

##############################

note_data['ha'] = 'left'
note_data['va'] = 'bottom'
note_data['r'] = 0
note_data.loc[note_data['grace']>1.2, 'ha'] = 'right'
note_data.loc[['Lee'], 'ha'] = 'right'
note_data.loc[['Christofferson', 'Monson', 'McConkie', 'Hunter', 'Maxwell', 'Uchtdorf',
               'Hinckley', 'Andersen', 'McKay', 'Faust', 'Smith', 'Benson'], 'va'] = 'top'
note_data.loc[['Hinckley', 'Maxwell', 'Benson'], 'ha'] = 'center'
note_data.loc[['Oaks', 'Cook'], 'r'] = 5
note_data.loc[['Ballard'], 'r'] = 10
titlestr = 'grace and works\n(per 1000 words)'
create_scatter_fig('grace', 'works', note_data, titlestr, 'grace_vs_works')






pres.save(args.output_dir + '/presentation.pptx')

